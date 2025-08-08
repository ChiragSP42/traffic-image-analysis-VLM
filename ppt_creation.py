from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
import os
import pandas as pd
import random

def add_slide(prs, title, left_text, image_path):
    # Use blank slide layout
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # Add heading at top center
    left = Inches(0.5)
    top = Inches(0.2)
    width = Inches(9)
    height = Inches(1)
    title_box = slide.shapes.add_textbox(left, top, width, height)
    title_frame = title_box.text_frame
    p = title_frame.paragraphs[0]
    run = p.add_run()
    run.text = title
    font = run.font
    font.name = 'Arial'
    font.size = Pt(32)
    font.bold = True
    font.color.rgb = RGBColor(0, 51, 102)
    p.alignment = 1  # 1=center

    # Add left-side text
    left_txt_left = Inches(0.5)
    left_txt_top = Inches(1.5)
    left_txt_width = Inches(4.25)
    left_txt_height = Inches(4)
    text_box = slide.shapes.add_textbox(left_txt_left, left_txt_top, left_txt_width, left_txt_height)
    text_frame = text_box.text_frame
    text_frame.text = left_text
    text_frame.paragraphs[0].font.size = Pt(18)
    text_frame.paragraphs[0].font.name = 'Calibri'
    text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)
    text_frame.word_wrap = True

    # Add right-side image (with error handling)
    img_left = Inches(5)
    img_top = Inches(1.5)
    img_width = Inches(4)
    if image_path and os.path.isfile(image_path):
        try:
            slide.shapes.add_picture(image_path, img_left, img_top, width=img_width)
        except Exception as e:
            print(f"Failed to insert image '{image_path}': {e}")
    else:
        print(f"Image not found: {image_path} -- Skipped adding image on slide '{title}'.")

def main():
    prs = Presentation()
    """
    1. First randomly select 20-30 images to be presented in the ppt.
    2. Extract these row indices and iterate one by one and pass the relevant information to ppt creation function.
    Format should be the following.
    slides_data = [
        {
            'title': '1999-2007-Black-Ford-F250-SuperDuty',
            'text': '4x4 badging on the side,\n Wooden rack/frame in truck bed,\n Tool box in truck bed,\n Extended cab configuration,\n Running boards/side steps,\n Visible license plate (partially visible),\n Standard chrome rear bumper.',
            'image': '1999-2007-Black-Ford-F250-SuperDuty.jpg'
        },
        {
            'title': '2000-2003-Gold-Nissan-Maxima',
            'text': 'Tinted windows, Five-spoke alloy wheels, Third brake light on trunk, Standard factory spoiler',
            'image': '2000-2003-Gold-Nissan-Maxima.jpg'
        },
        # Add more slides as needed
    ]
    """
    df = pd.read_csv("response_with_filenames.csv")
    count = 30
    rand_indices = random.sample(range(df.shape[0]), count)
    df = df.loc[rand_indices].copy()

    slides_data = list()
    for row in df.itertuples():
        title = row.File_name
        text = row.Response_identifiers
        image = f"{row.File_name}.jpg"
        slides_data.append({"title": title, "text": text, "image": image})

    for item in slides_data:
        add_slide(prs, item['title'], item['text'], os.path.join("Data", item['image']))

    prs.save('dataset_validation.pptx')
    print('Presentation created.')

if __name__ == '__main__':
    main()
